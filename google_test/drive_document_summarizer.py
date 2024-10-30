import logging
import sys
from typing import Dict, List, Optional
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from llama_index.core import Document, SummaryIndex, Settings
from llama_index.readers.file import PDFReader
from llama_index.readers.google import GoogleDocsReader
import json
import pandas as pd
import io
from googleapiclient.http import MediaIoBaseDownload
import tempfile
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from neo4j_test.user_store import UserStore
from pathlib import Path
from dotenv import load_dotenv
from llama_index.readers.smart_pdf_loader import SmartPDFLoader
import PyPDF2  
from docx import Document as DocxDocument
import docx2txt
from openpyxl import load_workbook
from pptx import Presentation
from .drive_stats import DriveStats

logging.basicConfig(stream=sys.stdout, level=logging.INFO)
logger = logging.getLogger(__name__)

load_dotenv()

class DriveDocumentSummarizer:
    # Default configuration
    DEFAULT_CONFIG = {
        'max_files_per_type': 20,  # Maximum files to process per type
        'max_total_files': 50,     # Maximum total files to process
        'test_mode': False,        # Test mode flag
        'test_file_limit': 10       # Number of files to process in test mode
    }
    
    MIME_TYPES = {
        'document': 'application/vnd.google-apps.document',
        'spreadsheet': 'application/vnd.google-apps.spreadsheet',
        'pdf': 'application/pdf',
        'presentation': 'application/vnd.google-apps.presentation',
        'word': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'excel': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'powerpoint': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'text': 'text/plain'
    }

    def __init__(self, credentials_dict: Dict = None, config: Dict = None):
        """Initialize with Google credentials and optional configuration"""
        # Merge provided config with defaults
        self.config = {**self.DEFAULT_CONFIG, **(config or {})}
        
        # Initialize credentials
        self.credentials = Credentials.from_authorized_user_info(
            credentials_dict,
            scopes=[
                "https://www.googleapis.com/auth/drive.metadata.readonly",
                "https://www.googleapis.com/auth/drive.readonly",
                "https://www.googleapis.com/auth/documents.readonly",
                "https://www.googleapis.com/auth/spreadsheets.readonly"
            ]
        )

        # Initialize services
        self.drive_service = build('drive', 'v3', credentials=self.credentials)
        self.docs_service = build('docs', 'v1', credentials=self.credentials)
        self.sheets_service = build('sheets', 'v4', credentials=self.credentials)
        self.docs_reader = GoogleDocsReader(credentials=self.credentials)
        
        # Initialize DriveStats
        self.drive_stats = DriveStats()

    def check_drive_stats(self, user_id: str) -> Dict:
        """Get Drive statistics before processing"""
        try:
            stats = self.drive_stats.get_file_count(user_id)
            logger.info(f"Drive Statistics for user {user_id}:")
            logger.info(f"Total Files: {stats['total_files']}")
            logger.info(f"Active Files: {stats['active_files']}")
            logger.info(f"Storage Used: {stats['total_size_human']}")
            return stats
        except Exception as e:
            logger.error(f"Error getting drive stats: {e}")
            return {}

    def list_files(self, user_id: str) -> Dict[str, List[Dict]]:
        """List files of each type in Google Drive with configured limits"""
        try:
            # First get drive statistics
            drive_stats = self.check_drive_stats(user_id)
            if not drive_stats:
                raise ValueError("Could not get drive statistics")

            files_by_type = {mime_type: [] for mime_type in self.MIME_TYPES.keys()}
            
            # Determine limits based on configuration and stats
            file_limit = (
                self.config['test_file_limit'] if self.config['test_mode']
                else min(self.config['max_files_per_type'], drive_stats['active_files'])
            )
            
            # Process each MIME type separately to ensure we get enough files of each type
            for doc_type, mime_type in self.MIME_TYPES.items():
                query = f"mimeType='{mime_type}' and trashed=false"
                
                try:
                    results = self.drive_service.files().list(
                        pageSize=file_limit,  # Increased from 10 to file_limit
                        fields="nextPageToken, files(id, name, mimeType)",
                        q=query,
                        orderBy="modifiedTime desc"
                    ).execute()
                    
                    items = results.get('files', [])
                    
                    if items:
                        # Take up to file_limit files of this type
                        for file in items[:file_limit]:
                            files_by_type[doc_type].append({
                                'id': file['id'],
                                'name': file['name'],
                                'mimeType': file['mimeType']
                            })
                            logger.info(f"Found {doc_type}: {file['name']}")
                
                except Exception as e:
                    logger.error(f"Error listing files of type {doc_type}: {e}")
                    continue
            
            # Log summary
            total_files = 0
            for doc_type, files in files_by_type.items():
                if files:
                    logger.info(f"Found {len(files)} {doc_type} files")
                    total_files += len(files)
            
            logger.info(f"Total files to process: {total_files}")
            return files_by_type

        except Exception as e:
            logger.error(f"Error listing files: {e}")
            raise

    def _summarize_document(self, doc_id: str) -> str:
        """Summarize a Google Doc"""
        try:
            logger.info(f"Attempting to read Google Doc: {doc_id}")
            
            # Get document content using Google Docs API
            doc = self.docs_service.documents().get(documentId=doc_id).execute()
            content = ""
            
            # Extract text from document
            for element in doc.get('body', {}).get('content', []):
                if 'paragraph' in element:
                    for para_element in element['paragraph']['elements']:
                        if 'textRun' in para_element:
                            content += para_element['textRun'].get('content', '')
            
            if not content.strip():
                logger.warning(f"No content found in Google Doc: {doc_id}")
                return "Unable to extract content from Google Doc"
            
            logger.info(f"Successfully loaded Google Doc. Content length: {len(content)}")
            
            # Create document and index
            document = Document(text=content)
            index = SummaryIndex.from_documents([document])
            query_engine = index.as_query_engine()
            
            summary_prompt = """Please provide a comprehensive summary of this document. Include:
            1. Main topics or themes
            2. Key points or findings
            3. Important details or conclusions
            Please structure the summary in a clear, organized manner."""
            
            response = query_engine.query(summary_prompt)
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing Google Doc {doc_id}: {str(e)}")
            return f"Error summarizing Google Doc: {str(e)}"

    def _summarize_spreadsheet(self, sheet_id: str) -> str:
        """Summarize a Google Spreadsheet"""
        try:
            logger.info(f"Attempting to read Google Spreadsheet: {sheet_id}")
            
            # Get spreadsheet metadata
            spreadsheet = self.sheets_service.spreadsheets().get(
                spreadsheetId=sheet_id
            ).execute()
            
            # Get file metadata to check size
            file_metadata = self.drive_service.files().get(
                fileId=sheet_id,
                fields='size'
            ).execute()
            
            file_size = int(file_metadata.get('size', '0'))
            size_limit = 14000  # Size threshold in bytes
            
            sheets_data = []
            
            # Process sheets based on size
            if file_size > size_limit:
                logger.info(f"Large spreadsheet detected ({file_size} bytes). Processing only first sheet.")
                # Get only the first sheet
                if spreadsheet.get('sheets'):
                    first_sheet = spreadsheet['sheets'][0]
                    sheet_title = first_sheet['properties']['title']
                    range_name = f"'{sheet_title}'"
                    
                    # Get sheet data
                    result = self.sheets_service.spreadsheets().values().get(
                        spreadsheetId=sheet_id,
                        range=range_name
                    ).execute()
                    
                    rows = result.get('values', [])
                    if rows:
                        sheet_content = [f"\nFirst Sheet: {sheet_title} (Note: Only showing first sheet due to large file size)"]
                        # Take only first 100 rows if there are more
                        rows = rows[:100]
                        for row in rows:
                            sheet_content.append('\t'.join(str(cell) for cell in row))
                        sheets_data.append('\n'.join(sheet_content))
            else:
                # Process all sheets for smaller files
                for sheet in spreadsheet.get('sheets', []):
                    sheet_title = sheet['properties']['title']
                    range_name = f"'{sheet_title}'"
                    
                    # Get sheet data
                    result = self.sheets_service.spreadsheets().values().get(
                        spreadsheetId=sheet_id,
                        range=range_name
                    ).execute()
                    
                    rows = result.get('values', [])
                    if rows:
                        sheet_content = [f"\nSheet: {sheet_title}"]
                        for row in rows:
                            sheet_content.append('\t'.join(str(cell) for cell in row))
                        sheets_data.append('\n'.join(sheet_content))
            
            if not sheets_data:
                logger.warning(f"No content found in spreadsheet: {sheet_id}")
                return "No content found in spreadsheet"
            
            full_content = '\n\n'.join(sheets_data)
            logger.info(f"Successfully loaded spreadsheet. Content length: {len(full_content)}")
            
            # Create document and index
            document = Document(text=full_content)
            index = SummaryIndex.from_documents([document])
            query_engine = index.as_query_engine()
            
            # Adjust summary prompt based on whether it's a full or partial summary
            if file_size > size_limit:
                summary_prompt = """Please analyze this spreadsheet data (first sheet only) and provide:
                1. Overview of the data structure in the first sheet
                2. Key patterns or trends visible in the available data
                3. Important insights from the first sheet
                Note: This is a partial analysis as only the first sheet was processed due to file size.
                Please structure the summary in a clear, organized manner."""
            else:
                summary_prompt = """Please analyze this spreadsheet data and provide:
                1. Overview of the data structure
                2. Key patterns or trends
                3. Important insights or findings
                Please structure the summary in a clear, organized manner."""
            
            response = query_engine.query(summary_prompt)
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing spreadsheet {sheet_id}: {str(e)}")
            return f"Error summarizing spreadsheet: {str(e)}"

    def _summarize_pdf(self, pdf_id: str) -> str:
        """Download and summarize a PDF file using PyPDF2"""
        temp_file = None
        try:
            # Get file metadata first
            file_metadata = self.drive_service.files().get(fileId=pdf_id).execute()
            
            # Create a temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
            
            # Download the file
            request = self.drive_service.files().get_media(fileId=pdf_id)
            downloader = MediaIoBaseDownload(temp_file, request)
            
            done = False
            while done is False:
                _, done = downloader.next_chunk()
            
            temp_file.close()
            
            # Read PDF with PyPDF2
            pdf_text = ""
            with open(temp_file.name, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract text from each page
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    pdf_text += f"\nPage {page_num + 1}:\n{page.extract_text()}"
            
            if not pdf_text.strip():
                return f"Unable to extract content from PDF: {file_metadata.get('name', 'Unknown file')}"
            
            # Create document and index
            document = Document(text=pdf_text)
            index = SummaryIndex.from_documents([document])
            query_engine = index.as_query_engine()
            
            # Create a more detailed summary prompt
            summary_prompt = """Please provide an overview of the content of this file. what is the purpose of the file? what is the main message? what are the key takeaways? """
            
            response = query_engine.query(summary_prompt)
            
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing PDF {pdf_id}: {str(e)}")
            return f"Error summarizing PDF: {str(e)}"
            
        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file.name):
                try:
                    os.unlink(temp_file.name)
                except Exception as e:
                    logger.error(f"Error cleaning up temporary file: {str(e)}")

    def _read_word(self, file_path: str) -> str:
        """Read content from Word document"""
        try:
            # First try with python-docx
            try:
                doc = DocxDocument(file_path)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            except:
                # If python-docx fails, try with docx2txt
                return docx2txt.process(file_path)
        except Exception as e:
            logger.error(f"Error reading Word document: {e}")
            return ""

    def _read_excel(self, file_path: str) -> str:
        """Read content from Excel file"""
        try:
            workbook = load_workbook(file_path)
            content = []
            
            for sheet in workbook.sheetnames:
                worksheet = workbook[sheet]
                content.append(f"\nSheet: {sheet}")
                
                for row in worksheet.iter_rows(values_only=True):
                    content.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
            
            return "\n".join(content)
        except Exception as e:
            logger.error(f"Error reading Excel file: {e}")
            return ""

    def _download_and_read_file(self, file_id: str, mime_type: str) -> str:
        """Download and read content from any file type"""
        temp_file = None
        try:
            # Create a temporary file
            suffix = self._get_file_extension(mime_type)
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
            
            # Download the file
            request = self.drive_service.files().get_media(fileId=file_id)
            downloader = MediaIoBaseDownload(temp_file, request)
            
            done = False
            while done is False:
                _, done = downloader.next_chunk()
            
            temp_file.close()
            
            # Read content based on file type
            if mime_type == self.MIME_TYPES['pdf']:
                return self._read_pdf(temp_file.name)
            elif mime_type == self.MIME_TYPES['word']:
                return self._read_word(temp_file.name)
            elif mime_type == self.MIME_TYPES['excel']:
                return self._read_excel(temp_file.name)
            elif mime_type == self.MIME_TYPES['text']:
                return self._read_text_file(temp_file.name)
            else:
                return f"File type {mime_type} not supported for direct reading"
                
        finally:
            if temp_file and os.path.exists(temp_file.name):
                os.unlink(temp_file.name)

    def _get_file_extension(self, mime_type: str) -> str:
        """Get file extension based on MIME type"""
        mime_to_ext = {
            self.MIME_TYPES['pdf']: '.pdf',
            self.MIME_TYPES['word']: '.docx',
            self.MIME_TYPES['excel']: '.xlsx',
            self.MIME_TYPES['powerpoint']: '.pptx',
            self.MIME_TYPES['text']: '.txt'
        }
        return mime_to_ext.get(mime_type, '')

    def _read_pdf(self, file_path: str) -> str:
        """Read content from PDF file"""
        pdf_text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                pdf_text += f"\nPage {page_num + 1}:\n{page.extract_text()}"
        return pdf_text

    def _read_text_file(self, file_path: str) -> str:
        """Read content from text-based files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # If UTF-8 fails, try with another common encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()

    def _read_presentation(self, file_path: str, max_slides: int = 3) -> str:
        """Read content from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            presentation_text = []
            
            # Add title slide if it exists
            if prs.slides:
                first_slide = prs.slides[0]
                if first_slide.shapes.title:
                    presentation_text.append(f"Title: {first_slide.shapes.title.text}")
            
            # Process up to max_slides
            for i, slide in enumerate(prs.slides):
                if i >= max_slides:
                    break
                    
                slide_text = [f"\nSlide {i+1}:"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                presentation_text.append("\n".join(slide_text))
            
            # Add note about limited slides
            if len(prs.slides) > max_slides:
                presentation_text.append(f"\nNote: Only showing first {max_slides} slides out of {len(prs.slides)} total slides.")
            
            return "\n\n".join(presentation_text)
            
        except Exception as e:
            logger.error(f"Error reading presentation: {e}")
            return ""

    def _summarize_presentation(self, presentation_id: str) -> str:
        """Download and summarize a presentation"""
        temp_file = None
        try:
            # Get file metadata
            file_metadata = self.drive_service.files().get(fileId=presentation_id).execute()
            
            # Check if it's a Google Slides presentation
            if file_metadata['mimeType'] == 'application/vnd.google-apps.presentation':
                # Export Google Slides to PPTX format
                request = self.drive_service.files().export_media(
                    fileId=presentation_id,
                    mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
            else:
                # For regular PPTX files, use normal download
                request = self.drive_service.files().get_media(fileId=presentation_id)
            
            # Create temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
            
            # Download file
            downloader = MediaIoBaseDownload(temp_file, request)
            
            done = False
            while done is False:
                _, done = downloader.next_chunk()
            
            temp_file.close()
            
            # Read presentation content
            presentation_text = self._read_presentation(temp_file.name)
            
            if not presentation_text.strip():
                return f"Unable to extract content from presentation: {file_metadata.get('name', 'Unknown file')}"
            
            # Create document and index
            document = Document(text=presentation_text)
            index = SummaryIndex.from_documents([document])
            query_engine = index.as_query_engine()
            
            summary_prompt = """Please analyze this presentation and provide:
            1. Main topic or purpose of the presentation
            2. Key points from the first few slides
            3. Important details or takeaways
            Please structure the summary in a clear, organized manner."""
            
            response = query_engine.query(summary_prompt)
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing presentation {presentation_id}: {str(e)}")
            return f"Error summarizing presentation: {str(e)}"
            
        finally:
            if temp_file and os.path.exists(temp_file.name):
                try:
                    os.unlink(temp_file.name)
                except Exception as e:
                    logger.error(f"Error cleaning up temporary file: {str(e)}")

    def _summarize_file(self, file_id: str, file_type: str, mime_type: str) -> str:
        """Summarize any type of file"""
        try:
            summary = ""
            error_msg = None
            
            try:
                if file_type == 'document':
                    summary = self._summarize_document(file_id)
                elif file_type == 'spreadsheet':
                    summary = self._summarize_spreadsheet(file_id)
                elif file_type == 'presentation':
                    summary = self._summarize_presentation(file_id)
                elif file_type in ['pdf', 'word', 'excel', 'text']:
                    content = self._download_and_read_file(file_id, mime_type)
                    if not content:
                        error_msg = "Unable to extract content from file"
                    else:
                        document = Document(text=content)
                        index = SummaryIndex.from_documents([document])
                        query_engine = index.as_query_engine()
                        
                        summary_prompt = f"""Please provide a comprehensive summary of this {file_type} file. Include:
                        1. Main purpose or topic of the document
                        2. Key points or findings
                        3. Important details or conclusions
                        4. Any significant data or metrics mentioned
                        
                        If the content appears to be unstructured or unclear, focus on identifying:
                        - The type of information present
                        - Any patterns or recurring themes
                        - The general context or domain
                        
                        Please structure the summary in a clear, organized manner."""
                        
                        response = query_engine.query(summary_prompt)
                        summary = str(response)
                else:
                    error_msg = f"File type {file_type} not supported for summarization"
                    
            except Exception as e:
                error_msg = f"Error during summarization: {str(e)}"
                logger.error(f"Error summarizing file {file_id}: {e}")
            
            # If summary is empty or too short, provide a status message
            if not summary or len(summary.strip()) < 50:
                if error_msg:
                    return f"Summary generation failed: {error_msg}"
                return "Unable to generate meaningful summary for this document"
                
            return summary
                
        except Exception as e:
            logger.error(f"Critical error summarizing file {file_id}: {e}")
            return f"Critical error in summarization process: {str(e)}"

    def summarize_all_files(self, user_id: str) -> Dict[str, List[Dict]]:
        """Summarize all files in Google Drive"""
        try:
            # Get files with user_id
            files_by_type = self.list_files(user_id)
            summaries = {doc_type: [] for doc_type in self.MIME_TYPES.keys()}
            
            for doc_type, files in files_by_type.items():
                for file in files:
                    try:
                        logger.info(f"Starting summarization of {doc_type}: {file['name']}")
                        mime_type = next((mt for dt, mt in self.MIME_TYPES.items() if dt == doc_type), None)
                        
                        # Get summary with error handling
                        try:
                            summary = self._summarize_file(file['id'], doc_type, mime_type)
                            logger.info(f"Successfully generated summary for {file['name']}")
                        except Exception as e:
                            logger.error(f"Error summarizing {file['name']}: {e}")
                            summary = f"Error generating summary: {str(e)}"
                        
                        summaries[doc_type].append({
                            'name': file['name'],
                            'id': file['id'],
                            'mime_type': mime_type,
                            'summary': summary,
                            'summary_status': 'SUCCESS' if len(summary) > 50 and not summary.startswith('Error') else 'FAILED'
                        })
                        
                    except Exception as e:
                        logger.error(f"Error processing file {file.get('name', 'unknown')}: {e}")
                        continue
            
            return summaries
            
        except Exception as e:
            logger.error(f"Error summarizing files: {e}")
            return {}

    def test_document_summary(self, doc_id: str) -> Dict:
        """Test summarization of a single document"""
        try:
            # Get document metadata
            file = self.drive_service.files().get(fileId=doc_id).execute()
            mime_type = file['mimeType']
            
            # Determine document type
            doc_type = next(
                (dt for dt, mt in self.MIME_TYPES.items() if mt == mime_type),
                'unknown'
            )
            
            # Generate summary
            summary = self._summarize_file(doc_id, doc_type, mime_type)
            
            return {
                'name': file['name'],
                'type': doc_type,
                'mime_type': mime_type,
                'summary': summary,
                'success': len(summary) > 50 and not summary.startswith('Error')
            }
            
        except Exception as e:
            logger.error(f"Error testing document summary: {e}")
            return {
                'error': str(e)
            }

def main():
    """Test the DriveDocumentSummarizer"""
    try:
        # Test configuration
        test_config = {
            'test_mode': True,
            'test_file_limit': 20,
            'max_total_files': 30
        }
        
        # Initialize UserStore to get credentials
        user_store = UserStore()
        test_user_id = "aman"
        
        # Get credentials
        credentials = user_store.get_google_credentials(test_user_id)
        if not credentials:
            logger.error("No valid credentials found")
            return
            
        # Initialize summarizer with credentials
        summarizer = DriveDocumentSummarizer(
            credentials_dict=credentials.to_json(),
            config=test_config
        )
        
        # Get drive statistics
        print("\nChecking Drive statistics...")
        stats = summarizer.check_drive_stats(test_user_id)
        print(f"Total files: {stats.get('total_files', 'N/A')}")
        print(f"Storage used: {stats.get('total_size_human', 'N/A')}")
        
        # Get and summarize files
        print("\nProcessing files...")
        files = summarizer.list_files(test_user_id)
        
        # Print results
        for doc_type, file_list in files.items():
            print(f"\n=== {doc_type.upper()} FILES ({len(file_list)}) ===")
            for file in file_list:
                print(f"- {file['name']}")
        
    except Exception as e:
        logger.error(f"Error in main: {e}")
        raise

if __name__ == "__main__":
    main()
