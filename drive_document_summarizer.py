import logging
import sys
from typing import Dict, List, Optional
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from llama_index.core import Document, SummaryIndex, Settings
from llama_index.readers.file import PDFReader  # Changed from llama_index.core.readers
from llama_index.readers.google import GoogleDocsReader
import json
import pandas as pd
import io
from googleapiclient.http import MediaIoBaseDownload
import tempfile
import os
from pathlib import Path
from dotenv import load_dotenv
from llama_index.readers.smart_pdf_loader import SmartPDFLoader
import PyPDF2  
from docx import Document as DocxDocument
import docx2txt
from openpyxl import load_workbook
from pptx import Presentation

logging.basicConfig(stream=sys.stdout, level=logging.INFO)
logger = logging.getLogger(__name__)

load_dotenv()

# Use API key from environment variable


class DriveDocumentSummarizer:
    # Update MIME_TYPES to use Google's native MIME types
    MIME_TYPES = {
        'document': 'application/vnd.google-apps.document', 
        'spreadsheet': 'application/vnd.google-apps.spreadsheet',  
        'pdf': 'application/pdf',
        'presentation': 'application/vnd.google-apps.presentation',
        'word': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'excel': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'powerpoint': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'text': 'text/plain'
    }

    def __init__(self, credentials_dict: Dict = None):
        """Initialize with Google credentials"""
        if credentials_dict is None:
            credentials_dict = self._get_credentials_from_db()
        
        # Update scopes to match exactly with google_doc.py
        self.credentials = Credentials.from_authorized_user_info(
            credentials_dict,
            scopes=[
                "https://www.googleapis.com/auth/drive.metadata.readonly",
                "https://www.googleapis.com/auth/drive.readonly",
                "https://www.googleapis.com/auth/documents.readonly",
                "https://www.googleapis.com/auth/spreadsheets.readonly"
            ]
        )
        
        if self.credentials.expired:
            self.credentials.refresh(Request())
            self._update_credentials_in_db(self.credentials.to_json())

        self.drive_service = build('drive', 'v3', credentials=self.credentials)
        self.docs_service = build('docs', 'v1', credentials=self.credentials)
        self.sheets_service = build('sheets', 'v4', credentials=self.credentials)
        self.docs_reader = GoogleDocsReader(credentials=self.credentials)

    def _get_credentials_from_db(self) -> Dict:
        """Get credentials from database (mock implementation)"""
        with open("token.json", "r") as f:
            return json.load(f)

    def _update_credentials_in_db(self, credentials_json: str):
        """Update credentials in database (mock implementation)"""
        with open("token.json", "w") as f:
            f.write(credentials_json)

    def list_files(self) -> Dict[str, List[Dict]]:
        """List first 10 files of each type in Google Drive"""
        try:
            files_by_type = {mime_type: [] for mime_type in self.MIME_TYPES.keys()}
            
            # Create query for Google's native types
            query_parts = []
            for doc_type, mime_type in self.MIME_TYPES.items():
                query_parts.append(f"mimeType='{mime_type}'")
            query = " or ".join(query_parts)
            
            logger.info("Querying Google Drive...")
            results = self.drive_service.files().list(
                pageSize=20,  # Keep this higher to get enough files to filter
                fields="nextPageToken, files(id, name, mimeType)",
                q=query,
                orderBy="modifiedTime desc"  # Get most recently modified files first
            ).execute()
            
            items = results.get('files', [])
            
            if not items:
                logger.warning("No files found.")
                return files_by_type
            
            # Categorize files by type, limiting to 10 per type change this later as per cron job --- to be discussed with Maheedhar
            for file in items:
                for doc_type, mime_type in self.MIME_TYPES.items():
                    if file['mimeType'] == mime_type:
                        # Only add if we haven't reached 10 files for this type
                        if len(files_by_type[doc_type]) < 10:
                            files_by_type[doc_type].append({
                                'id': file['id'],
                                'name': file['name'],
                                'mimeType': file['mimeType']
                            })
                            logger.info(f"Found {doc_type}: {file['name']}")
                        break  # Stop checking other mime types once we've found a match
            
            # Log summary of files found
            for doc_type, files in files_by_type.items():
                if files:
                    logger.info(f"Found {len(files)} {doc_type} files")
            
            return files_by_type

        except Exception as e:
            logger.error(f"Error listing files: {e}")
            raise

    def _summarize_document(self, doc_id: str) -> str:
        """Summarize a Google Doc"""
        try:
            logger.info(f"Attempting to read Google Doc: {doc_id}")
            documents = self.docs_reader.load_data(document_ids=[doc_id])
            
            if not documents:
                logger.warning(f"No content found in Google Doc: {doc_id}")
                return "Unable to extract content from Google Doc"
            
            logger.info(f"Successfully loaded Google Doc. Content length: {len(documents[0].text)}")
            
            # Create index and generate summary
            index = SummaryIndex.from_documents(documents)
            query_engine = index.as_query_engine()
            
            summary_prompt = """Please provide a comprehensive summary of this document. Include:
            1. Main topics or themes
            2. Key points or findings
            3. Important details or conclusions
            Please structure the summary in a clear, organized manner."""
            
            response = query_engine.query(summary_prompt)
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing Google Doc {doc_id}: {str(e)}")
            return f"Error summarizing Google Doc: {str(e)}"

    def _summarize_spreadsheet(self, sheet_id: str) -> str:
        """Summarize a Google Spreadsheet"""
        try:
            logger.info(f"Attempting to read Google Spreadsheet: {sheet_id}")
            
            # Get spreadsheet metadata
            spreadsheet = self.sheets_service.spreadsheets().get(
                spreadsheetId=sheet_id
            ).execute()
            
            sheets_data = []
            
            # Process each sheet
            for sheet in spreadsheet.get('sheets', []):
                sheet_title = sheet['properties']['title']
                range_name = f"'{sheet_title}'"
                
                # Get sheet data
                result = self.sheets_service.spreadsheets().values().get(
                    spreadsheetId=sheet_id,
                    range=range_name
                ).execute()
                
                rows = result.get('values', [])
                if rows:
                    sheet_content = [f"\nSheet: {sheet_title}"]
                    for row in rows:
                        sheet_content.append('\t'.join(str(cell) for cell in row))
                    sheets_data.append('\n'.join(sheet_content))
            
            if not sheets_data:
                logger.warning(f"No content found in spreadsheet: {sheet_id}")
                return "No content found in spreadsheet"
            
            full_content = '\n\n'.join(sheets_data)
            logger.info(f"Successfully loaded spreadsheet. Content length: {len(full_content)}")
            
            # Create document and index
            document = Document(text=full_content)
            index = SummaryIndex.from_documents([document])
            query_engine = index.as_query_engine()
            
            summary_prompt = """Please analyze this spreadsheet data and provide:
            1. Overview of the data structure
            2. Key patterns or trends
            3. Important insights or findings
            Please structure the summary in a clear, organized manner."""
            
            response = query_engine.query(summary_prompt)
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing spreadsheet {sheet_id}: {str(e)}")
            return f"Error summarizing spreadsheet: {str(e)}"

    def _summarize_pdf(self, pdf_id: str) -> str:
        """Download and summarize a PDF file using PyPDF2"""
        temp_file = None
        try:
            # Get file metadata first
            file_metadata = self.drive_service.files().get(fileId=pdf_id).execute()
            
            # Create a temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
            
            # Download the file
            request = self.drive_service.files().get_media(fileId=pdf_id)
            downloader = MediaIoBaseDownload(temp_file, request)
            
            done = False
            while done is False:
                _, done = downloader.next_chunk()
            
            temp_file.close()
            
            # Read PDF with PyPDF2
            pdf_text = ""
            with open(temp_file.name, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Extract text from each page
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    pdf_text += f"\nPage {page_num + 1}:\n{page.extract_text()}"
            
            if not pdf_text.strip():
                return f"Unable to extract content from PDF: {file_metadata.get('name', 'Unknown file')}"
            
            # Create document and index
            document = Document(text=pdf_text)
            index = SummaryIndex.from_documents([document])
            query_engine = index.as_query_engine()
            
            # Create a more detailed summary prompt
            summary_prompt = """Please provide an overview of the content of this file. what is the purpose of the file? what is the main message? what are the key takeaways? """
            
            response = query_engine.query(summary_prompt)
            
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing PDF {pdf_id}: {str(e)}")
            return f"Error summarizing PDF: {str(e)}"
            
        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file.name):
                try:
                    os.unlink(temp_file.name)
                except Exception as e:
                    logger.error(f"Error cleaning up temporary file: {str(e)}")

    def _read_word(self, file_path: str) -> str:
        """Read content from Word document"""
        try:
            # First try with python-docx
            try:
                doc = DocxDocument(file_path)
                return '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            except:
                # If python-docx fails, try with docx2txt
                return docx2txt.process(file_path)
        except Exception as e:
            logger.error(f"Error reading Word document: {e}")
            return ""

    def _read_excel(self, file_path: str) -> str:
        """Read content from Excel file"""
        try:
            workbook = load_workbook(file_path)
            content = []
            
            for sheet in workbook.sheetnames:
                worksheet = workbook[sheet]
                content.append(f"\nSheet: {sheet}")
                
                for row in worksheet.iter_rows(values_only=True):
                    content.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
            
            return "\n".join(content)
        except Exception as e:
            logger.error(f"Error reading Excel file: {e}")
            return ""

    def _download_and_read_file(self, file_id: str, mime_type: str) -> str:
        """Download and read content from any file type"""
        temp_file = None
        try:
            # Create a temporary file
            suffix = self._get_file_extension(mime_type)
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
            
            # Download the file
            request = self.drive_service.files().get_media(fileId=file_id)
            downloader = MediaIoBaseDownload(temp_file, request)
            
            done = False
            while done is False:
                _, done = downloader.next_chunk()
            
            temp_file.close()
            
            # Read content based on file type
            if mime_type == self.MIME_TYPES['pdf']:
                return self._read_pdf(temp_file.name)
            elif mime_type == self.MIME_TYPES['word']:
                return self._read_word(temp_file.name)
            elif mime_type == self.MIME_TYPES['excel']:
                return self._read_excel(temp_file.name)
            elif mime_type == self.MIME_TYPES['text']:
                return self._read_text_file(temp_file.name)
            else:
                return f"File type {mime_type} not supported for direct reading"
                
        finally:
            if temp_file and os.path.exists(temp_file.name):
                os.unlink(temp_file.name)

    def _get_file_extension(self, mime_type: str) -> str:
        """Get file extension based on MIME type"""
        mime_to_ext = {
            self.MIME_TYPES['pdf']: '.pdf',
            self.MIME_TYPES['word']: '.docx',
            self.MIME_TYPES['excel']: '.xlsx',
            self.MIME_TYPES['powerpoint']: '.pptx',
            self.MIME_TYPES['text']: '.txt'
        }
        return mime_to_ext.get(mime_type, '')

    def _read_pdf(self, file_path: str) -> str:
        """Read content from PDF file"""
        pdf_text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                pdf_text += f"\nPage {page_num + 1}:\n{page.extract_text()}"
        return pdf_text

    def _read_text_file(self, file_path: str) -> str:
        """Read content from text-based files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # If UTF-8 fails, try with another common encoding
            with open(file_path, 'r', encoding='latin-1') as file:
                return file.read()

    def _read_presentation(self, file_path: str, max_slides: int = 3) -> str:
        """Read content from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            presentation_text = []
            
            # Add title slide if it exists
            if prs.slides:
                first_slide = prs.slides[0]
                if first_slide.shapes.title:
                    presentation_text.append(f"Title: {first_slide.shapes.title.text}")
            
            # Process up to max_slides
            for i, slide in enumerate(prs.slides):
                if i >= max_slides:
                    break
                    
                slide_text = [f"\nSlide {i+1}:"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                presentation_text.append("\n".join(slide_text))
            
            # Add note about limited slides
            if len(prs.slides) > max_slides:
                presentation_text.append(f"\nNote: Only showing first {max_slides} slides out of {len(prs.slides)} total slides.")
            
            return "\n\n".join(presentation_text)
            
        except Exception as e:
            logger.error(f"Error reading presentation: {e}")
            return ""

    def _summarize_presentation(self, presentation_id: str) -> str:
        """Download and summarize a presentation"""
        temp_file = None
        try:
            # Get file metadata
            file_metadata = self.drive_service.files().get(fileId=presentation_id).execute()
            
            # Check if it's a Google Slides presentation
            if file_metadata['mimeType'] == 'application/vnd.google-apps.presentation':
                # Export Google Slides to PPTX format
                request = self.drive_service.files().export_media(
                    fileId=presentation_id,
                    mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
            else:
                # For regular PPTX files, use normal download
                request = self.drive_service.files().get_media(fileId=presentation_id)
            
            # Create temporary file
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
            
            # Download file
            downloader = MediaIoBaseDownload(temp_file, request)
            
            done = False
            while done is False:
                _, done = downloader.next_chunk()
            
            temp_file.close()
            
            # Read presentation content
            presentation_text = self._read_presentation(temp_file.name)
            
            if not presentation_text.strip():
                return f"Unable to extract content from presentation: {file_metadata.get('name', 'Unknown file')}"
            
            # Create document and index
            document = Document(text=presentation_text)
            index = SummaryIndex.from_documents([document])
            query_engine = index.as_query_engine()
            
            summary_prompt = """Please analyze this presentation and provide:
            1. Main topic or purpose of the presentation
            2. Key points from the first few slides
            3. Important details or takeaways
            Please structure the summary in a clear, organized manner."""
            
            response = query_engine.query(summary_prompt)
            return str(response)
            
        except Exception as e:
            logger.error(f"Error summarizing presentation {presentation_id}: {str(e)}")
            return f"Error summarizing presentation: {str(e)}"
            
        finally:
            if temp_file and os.path.exists(temp_file.name):
                try:
                    os.unlink(temp_file.name)
                except Exception as e:
                    logger.error(f"Error cleaning up temporary file: {str(e)}")

    def _summarize_file(self, file_id: str, file_type: str, mime_type: str) -> str:
        """Summarize any type of file"""
        try:
            if file_type == 'document':
                return self._summarize_document(file_id)
            elif file_type == 'spreadsheet':
                return self._summarize_spreadsheet(file_id)
            elif file_type == 'presentation':
                return self._summarize_presentation(file_id)
            elif file_type in ['pdf', 'word', 'excel', 'text']:
                # Download and read content
                content = self._download_and_read_file(file_id, mime_type)
                if not content:
                    return "Unable to extract content from file"
                
                # Create document and index
                document = Document(text=content)
                index = SummaryIndex.from_documents([document])
                query_engine = index.as_query_engine()
                
                summary_prompt = """Please provide an overview of the content of this file. What is the purpose of the file? What is the main message? What are the key takeaways? Please structure the response in a clear, organized manner."""
                
                response = query_engine.query(summary_prompt)
                return str(response)
            else:
                return "File type not supported for summarization"
                
        except Exception as e:
            logger.error(f"Error summarizing file {file_id}: {e}")
            return f"Error summarizing file: {str(e)}"

    def summarize_all_files(self) -> Dict[str, List[Dict]]:
        """Summarize all files in Google Drive"""
        files_by_type = self.list_files()
        summaries = {doc_type: [] for doc_type in self.MIME_TYPES.keys()}
        
        for doc_type, files in files_by_type.items():
            for file in files:
                logger.info(f"Summarizing {doc_type}: {file['name']}")
                mime_type = next((mt for dt, mt in self.MIME_TYPES.items() if dt == doc_type), None)
                summary = self._summarize_file(file['id'], doc_type, mime_type)
                
                summaries[doc_type].append({
                    'name': file['name'],
                    'id': file['id'],
                    'summary': summary
                })
        
        return summaries

def main():
    try:
        summarizer = DriveDocumentSummarizer()
        
        # Get and summarize all files
        summaries = summarizer.summarize_all_files()
        
        # Print results
        for doc_type, files in summaries.items():
            print(f"\n=== {doc_type.upper()} FILES ===")
            for file in files:
                print(f"\nFile: {file['name']}")
                print(f"ID: {file['id']}")
                print("Summary:")
                print(file['summary'])
                print("-" * 80)
        
    except Exception as e:
        logger.error(f"Error in main: {e}")
        raise

if __name__ == "__main__":
    main()
